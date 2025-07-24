from pptx import Presentation
from PIL import Image
import io

# 压缩图像
def compress_image(image_path, output_path, quality=20):
    with Image.open(image_path) as img:
        img.save(output_path, "JPEG", quality=quality)

# 处理PPT文件
def optimize_ppt(input_ppt, output_ppt):
    prs = Presentation(input_ppt)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # 图片类型
                image_stream = io.BytesIO(shape.image.blob)
                img = Image.open(image_stream)

                # 生成压缩后的图片
                compressed_image_path = 'temp_compressed.jpg'
                compress_image(image_stream, compressed_image_path)

                # 替换图片
                with open(compressed_image_path, 'rb') as f:
                    shape.image = f.read()

    prs.save(output_ppt)

# 示例
optimize_ppt('input.pptx', 'output.pptx')
